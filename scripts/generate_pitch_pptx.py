#!/usr/bin/env python3
"""从 docs/intro/CoAgent-路演PPT.html 生成可编辑 PowerPoint 文件。"""

import base64
import re
import sys
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
HTML = ROOT / "docs/intro/CoAgent-路演PPT.html"
OUT = ROOT / "docs/intro/CoAgent-路演PPT.pptx"

BG = RGBColor(0x04, 0x09, 0x1F)
PANEL = RGBColor(0x0C, 0x17, 0x38)
CYAN = RGBColor(0x00, 0xE0, 0xFF)
BLUE_L = RGBColor(0x6E, 0xC5, 0xFF)
MUTED = RGBColor(0x6F, 0x86, 0xB3)
TEXT = RGBColor(0xEA, 0xF1, 0xFF)
DIM = RGBColor(0xB8, 0xC9, 0xE6)
GOLD = RGBColor(0xFF, 0xD1, 0x66)
LINE = RGBColor(0x21, 0x30, 0x5C)
WARN = RGBColor(0xFF, 0xB0, 0x20)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text="", size=14, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    p.alignment = align
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "PingFang SC"
    return tb


def add_footer(slide, page, total=8):
    add_textbox(slide, Inches(0.6), Inches(7.05), Inches(5), Inches(0.3), "小宿科技环球黑客松 · 北京站", 9, MUTED)
    add_textbox(slide, Inches(11.5), Inches(7.05), Inches(1.2), Inches(0.3), f"{page:02d} / {total:02d}", 9, BLUE_L, True, PP_ALIGN.RIGHT)


def add_seclabel(slide, en, cn):
    tb = slide.shapes.add_textbox(Inches(0.65), Inches(0.45), Inches(10), Inches(0.35))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = en.upper() + "  "
    r1.font.size = Pt(11)
    r1.font.color.rgb = CYAN
    r1.font.bold = True
    r1.font.name = "PingFang SC"
    r2 = p.add_run()
    r2.text = "· " + cn
    r2.font.size = Pt(11)
    r2.font.color.rgb = MUTED
    r2.font.name = "PingFang SC"


def add_title(slide, title):
    add_textbox(slide, Inches(0.65), Inches(0.85), Inches(12), Inches(0.7), title, 28, WHITE, True)


def extract_cover_images(html: str, dest: Path) -> list[Path]:
    pattern = re.compile(r'<img\s+src="(data:image/[^;]+;base64,[^"]+)"')
    saved = []
    for i, data_uri in enumerate(pattern.findall(html)[:6]):
        header, b64 = data_uri.split(",", 1)
        ext = "png" if "png" in header else "jpeg"
        path = dest / f"cover_{i}.{ext}"
        path.write_bytes(base64.b64decode(b64))
        saved.append(path)
    return saved


def build_presentation() -> Presentation:
    html = HTML.read_text(encoding="utf-8")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    with tempfile.TemporaryDirectory() as tmp:
        imgs = extract_cover_images(html, Path(tmp))

        # Slide 1 — Cover
        s1 = prs.slides.add_slide(blank)
        add_bg(s1)
        if imgs:
            try:
                s1.shapes.add_picture(str(imgs[0]), Inches(0), Inches(0), width=Inches(13.333))
            except Exception:
                pass
        for img, x in zip(imgs[1:3], [Inches(0.5), Inches(11.8)]):
            try:
                s1.shapes.add_picture(str(img), x, Inches(0.35), height=Inches(0.55))
            except Exception:
                pass
        add_textbox(s1, Inches(0.65), Inches(1.8), Inches(8), Inches(0.35), "Beyond Prompt · Agents in Action", 12, CYAN, True)
        add_textbox(s1, Inches(0.65), Inches(2.2), Inches(10), Inches(1.2), "CoAgent", 72, CYAN, True)
        add_textbox(s1, Inches(0.65), Inches(3.4), Inches(10), Inches(0.8), "为 Agent 上生产护航：诊断「错在哪里」，研判「怎么处理」，协同「帮你处理」", 16, DIM)
        add_textbox(s1, Inches(0.65), Inches(4.5), Inches(2), Inches(0.25), "TRACK ·", 9, MUTED, True)
        add_textbox(s1, Inches(0.65), Inches(4.75), Inches(5), Inches(0.4), "赛道三 · ToB 场景 AI Agent（开放赛道）", 14)
        add_textbox(s1, Inches(0.65), Inches(5.3), Inches(1.5), Inches(0.25), "TEAM ·", 9, MUTED, True)
        add_textbox(s1, Inches(0.65), Inches(5.55), Inches(4), Inches(0.35), "Beyond", 14)
        add_textbox(s1, Inches(3.5), Inches(5.3), Inches(1.5), Inches(0.25), "MEMBERS ·", 9, MUTED, True)
        add_textbox(s1, Inches(3.5), Inches(5.55), Inches(4), Inches(0.35), "路睿 · 葛丽莎", 14)
        add_textbox(s1, Inches(6.5), Inches(5.3), Inches(1.5), Inches(0.25), "DATE ·", 9, MUTED, True)
        add_textbox(s1, Inches(6.5), Inches(5.55), Inches(5), Inches(0.35), "2026.06.28 · 北京·微软亚太研发集团总部", 14)
        add_textbox(s1, Inches(10.5), Inches(0.65), Inches(2.2), Inches(0.25), "PITCH DURATION", 9, MUTED, True, PP_ALIGN.RIGHT)
        add_textbox(s1, Inches(10.5), Inches(0.9), Inches(2.2), Inches(0.6), "04:00", 28, CYAN, True, PP_ALIGN.RIGHT)
        add_textbox(s1, Inches(10.5), Inches(1.45), Inches(2.2), Inches(0.3), "+ 2min QA", 10, MUTED, False, PP_ALIGN.RIGHT)
        add_footer(s1, 1)

        # Slide 2 — Problem
        s2 = prs.slides.add_slide(blank)
        add_bg(s2)
        add_seclabel(s2, "THE PROBLEM", "痛点与机会")
        add_title(s2, "Agent 加速进入企业，治理断层正在放大")
        add_textbox(s2, Inches(0.65), Inches(1.55), Inches(11), Inches(0.45), "采用率快速上升，但治理能力没有同步建立，最终把生产问题推向项目失败。", 13, DIM)
        stats = [
            ("ADOPTION · AGENT 采用率", "62%", "企业已试验 Agent", "但近 2/3 尚未进入规模化部署。", "McKinsey · 2025", False),
            ("GOVERNANCE GAP · 治理缺口", "21%", "具备成熟治理能力", "大量 Agent 已上线，却缺少生产治理体系。", "Deloitte · 2026", False),
            ("FAILURE RISK · 失败风险", "40%+", "Agent 项目可能被取消", "到 2027 年底，成本、价值与风险失控将成为主因。", "Gartner · 2027", True),
        ]
        for i, (eyebrow, num, h3, desc, source, is_risk) in enumerate(stats):
            x = Inches(0.65 + i * 4.1)
            fc = RGBColor(0x2A, 0x1D, 0x07) if is_risk else PANEL
            lc = RGBColor(0x5A, 0x44, 0x13) if is_risk else LINE
            add_rect(s2, x, Inches(2.15), Inches(3.85), Inches(1.65), fc, lc)
            add_textbox(s2, x + Inches(0.15), Inches(2.25), Inches(3.5), Inches(0.2), eyebrow, 7, MUTED, True)
            add_textbox(s2, x + Inches(0.15), Inches(2.45), Inches(3.5), Inches(0.55), num, 36, GOLD if is_risk else CYAN, True)
            add_textbox(s2, x + Inches(0.15), Inches(3.0), Inches(3.5), Inches(0.3), h3, 13, WHITE, True)
            add_textbox(s2, x + Inches(0.15), Inches(3.3), Inches(3.5), Inches(0.35), desc, 9, MUTED)
            add_textbox(s2, x + Inches(0.15), Inches(3.65), Inches(3.5), Inches(0.2), source, 7, RGBColor(0x46, 0x59, 0x82))
        add_rect(s2, Inches(0.65), Inches(4.0), Inches(12), Inches(0.55), RGBColor(0x09, 0x14, 0x31), LINE)
        for i, p in enumerate(["质量问题\n看不见", "运行异常\n找不到", "成本失控\n止不住", "处置风险\n不敢动"]):
            add_textbox(s2, Inches(0.65 + i * 3.0), Inches(4.05), Inches(2.9), Inches(0.5), p, 11, DIM, False, PP_ALIGN.CENTER)
        add_textbox(s2, Inches(0.65), Inches(4.75), Inches(8), Inches(0.25), "MARKET CAPABILITIES · 当前市场已有能力", 9, CYAN, True)
        add_textbox(s2, Inches(8), Inches(4.75), Inches(4.5), Inches(0.25), "能力存在，但分散在不同工具和工作流中", 9, MUTED, False, PP_ALIGN.RIGHT)
        tools = [
            ("① 看见", "Trace · 指标 · Eval · 告警", "LangSmith · Datadog"),
            ("② 诊断", "根因调查 · 关联证据 · 修复建议", "Galileo · Rootly"),
            ("③ 控制", "Guardrail · 策略阻断 · Runbook", "Cisco · Snyk"),
            ("④ 协同", "IM 卡片 · 审批 · 升级", "PagerDuty · Rootly"),
        ]
        for i, (h, p, prod) in enumerate(tools):
            x = Inches(0.65 + i * 3.05)
            add_rect(s2, x, Inches(5.05), Inches(2.9), Inches(0.95), PANEL, LINE)
            add_textbox(s2, x + Inches(0.1), Inches(5.1), Inches(2.7), Inches(0.25), h, 11, WHITE, True)
            add_textbox(s2, x + Inches(0.1), Inches(5.35), Inches(2.7), Inches(0.3), p, 8.5, DIM)
            add_textbox(s2, x + Inches(0.1), Inches(5.65), Inches(2.7), Inches(0.2), prod, 7.5, MUTED)
        add_rect(s2, Inches(0.65), Inches(6.15), Inches(12), Inches(0.45), RGBColor(0x00, 0x22, 0x33), CYAN)
        add_textbox(s2, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.4), "缺口：不是「没有能力」，而是缺少低集成成本、Agent 原生、面向事故责任人的一体化处置闭环。", 11, DIM)
        add_footer(s2, 2)

        # Slide 3 — Solution
        s3 = prs.slides.add_slide(blank)
        add_bg(s3)
        add_seclabel(s3, "THE SOLUTION", "我们的 Agent")
        add_title(s3, "我们的解法 · 一个能干活的 Agent")
        add_textbox(s3, Inches(0.65), Inches(1.55), Inches(11.5), Inches(0.55), "我们做了一个「Agent 事故指挥官」，让 FDE / 值班 / 客户 IT 用一张卡片完成「从诊断到处置生产事故」。", 14, DIM)
        caps = [
            ("01 CAPABILITY", "👁 看见问题", "接入运行事件与指标（POST /events）、幂等去重，把运行失败 / 质量异常 / 成本失控识别为「事故」。"),
            ("02 CAPABILITY", "🔍 诊断问题", "汇总日志/指标/配置/手册 → 根因假设 + 证据链；把握度评分给出「该信什么、有多大把握」。"),
            ("03 CAPABILITY", "🛠 处理问题", "处置合同（🟢 可执行 / 🟡 需确认 / 🔴 必须升级）+ 飞书协同 + 受控动作，帮你安全动手。"),
        ]
        for i, (kick, h, p) in enumerate(caps):
            x = Inches(0.65 + i * 4.1)
            add_rect(s3, x, Inches(2.3), Inches(3.85), Inches(2.5), PANEL, LINE)
            add_textbox(s3, x + Inches(0.2), Inches(2.4), Inches(3.5), Inches(0.2), kick, 8, MUTED, True)
            add_textbox(s3, x + Inches(0.2), Inches(2.65), Inches(3.5), Inches(0.35), h, 16, WHITE, True)
            add_textbox(s3, x + Inches(0.2), Inches(3.05), Inches(3.45), Inches(1.5), p, 11, DIM)
        add_rect(s3, Inches(0.65), Inches(5.1), Inches(12), Inches(0.55), RGBColor(0x00, 0x22, 0x33), CYAN)
        add_textbox(s3, Inches(0.8), Inches(5.18), Inches(11.5), Inches(0.45), "▸ 差异：别人把「看见、诊断、处理」拆在不同工具、还停在屏幕上；CoAgent 收进一个闭环，帮你动手。", 12, DIM)
        add_footer(s3, 3)

        # Slide 4 — Demo
        s4 = prs.slides.add_slide(blank)
        add_bg(s4)
        add_seclabel(s4, "PRODUCT DEMO", "看它怎么干活")
        add_title(s4, "Demo · 看它怎么干活")
        demo = ROOT / "docs/demos/coagent-demo.gif"
        add_rect(s4, Inches(0.65), Inches(1.55), Inches(6.8), Inches(3.85), RGBColor(0x01, 0x04, 0x0D), RGBColor(0x2C, 0x61, 0xA8))
        if demo.exists():
            s4.shapes.add_picture(str(demo), Inches(0.65), Inches(1.55), width=Inches(6.8))
        add_textbox(s4, Inches(0.85), Inches(1.7), Inches(2.5), Inches(0.25), "15 SEC · LIVE PRODUCT", 8, CYAN, True)
        steps = [
            ("01 →", "触发事件", "控制条选择场景，或由真实 Agent 回调 POST /events。"),
            ("02 →", "看见 + 诊断", "汇总运行信号，生成根因、证据链与把握度评分。"),
            ("03 →", "帮你处理", "输出风险分级处置合同，并通过飞书协同确认。"),
            ("04 →", "复盘沉淀", "验证处置结果，沉淀审计记录与反馈飞轮。"),
        ]
        for i, (s, h, p) in enumerate(steps):
            y = Inches(1.55 + i * 0.92)
            add_rect(s4, Inches(7.7), y, Inches(5.0), Inches(0.82), PANEL, LINE)
            add_textbox(s4, Inches(7.85), y + Inches(0.05), Inches(0.6), Inches(0.25), s, 9, CYAN, True)
            add_textbox(s4, Inches(8.4), y + Inches(0.05), Inches(4.1), Inches(0.25), h, 11, WHITE, True)
            add_textbox(s4, Inches(8.4), y + Inches(0.32), Inches(4.1), Inches(0.4), p, 9, DIM)
        add_textbox(s4, Inches(0.65), Inches(5.55), Inches(12), Inches(0.35), "同一套引擎，三种事故三种处理：  S1 87 🟢 重试    S2 72 🟡 禁试    S3 58 🔴 升级", 12, DIM)
        add_footer(s4, 4)

        # Slide 5 — Architecture
        s5 = prs.slides.add_slide(blank)
        add_bg(s5)
        add_seclabel(s5, "ARCHITECTURE", "技术实现")
        add_title(s5, "技术实现 · 异步可验证的处置闭环")
        arch = ROOT / "docs/diagrams/coagent-architecture.png"
        if arch.exists():
            s5.shapes.add_picture(str(arch), Inches(0.65), Inches(1.45), width=Inches(12.0))
        add_footer(s5, 5)

        # Slide 6 — Business
        s6 = prs.slides.add_slide(blank)
        add_bg(s6)
        add_seclabel(s6, "BUSINESS VALUE", "谁会为它买单")
        add_title(s6, "商业价值 · 谁会为它买单")
        blocks = [
            ("Target Customers · 客户", ["Agent 交付团队：FDE / 解决方案商，少人维护多客户", "企业 AI 平台 / SRE：保线上稳定、控成本", "客户 IT / 业务负责人：看懂依据、完成审批"]),
            ("Core Scenarios · 核心场景", ["S1 限流止损（429 → 一键重试）", "S2 空检索拦截（禁止无效重试）", "S3 成本超预算（升级审批）"]),
        ]
        for i, (head, items) in enumerate(blocks):
            x = Inches(0.65 + i * 6.2)
            add_rect(s6, x, Inches(1.55), Inches(5.9), Inches(2.2), PANEL, LINE)
            add_textbox(s6, x + Inches(0.15), Inches(1.65), Inches(5.6), Inches(0.25), head, 9, CYAN, True)
            for j, item in enumerate(items):
                add_textbox(s6, x + Inches(0.15), Inches(2.0 + j * 0.45), Inches(5.6), Inches(0.4), "▸ " + item, 11, DIM)
        add_textbox(s6, Inches(0.65), Inches(3.95), Inches(4), Inches(0.25), "Quantified Impact · 量化收益", 9, CYAN, True)
        for i, (warn, num, label) in enumerate([("数字待估 ⚠️", "N×", "诊断/MTTR 提效"), ("", "-XX%", "止损/超支下降"), ("", "¥X", "单客户 ARR")]):
            x = Inches(0.65 + i * 2.5)
            if warn:
                add_textbox(s6, x, Inches(4.2), Inches(2.2), Inches(0.2), warn, 8, WARN)
            add_textbox(s6, x, Inches(4.45), Inches(2.2), Inches(0.45), num, 24, CYAN, True)
            add_textbox(s6, x, Inches(4.95), Inches(2.2), Inches(0.25), label, 8, MUTED)
        add_rect(s6, Inches(0.65), Inches(5.35), Inches(12), Inches(1.35), PANEL, LINE)
        add_textbox(s6, Inches(0.8), Inches(5.45), Inches(4), Inches(0.25), "Business Model · 商业模式", 9, CYAN, True)
        add_textbox(s6, Inches(0.8), Inches(5.72), Inches(3), Inches(0.2), "假设待定稿 ⚠️", 8, WARN)
        for j, item in enumerate(["收费：按接入 Agent 数 / 席位订阅", "GTM：先切 FDE 交付团队", "护城河：处置合同机制 + 结果验证飞轮"]):
            add_textbox(s6, Inches(0.8 + j * 4.0), Inches(6.05), Inches(3.8), Inches(0.35), "▸ " + item, 11, DIM)
        add_footer(s6, 6)

        # Slide 7 — 48 Hours
        s7 = prs.slides.add_slide(blank)
        add_bg(s7)
        add_seclabel(s7, "48 HOURS", "我们做到了什么")
        add_title(s7, "48 小时 · 我们做到了什么")
        timeline = [
            ("10:00 · DAY 1", "组队 & 选题", "锁定「Agent 事故处置 / SRE war room」方向，定下三场景。"),
            ("18:00–24:00 · DAY 1", "核心链路打通", "P0 baseline 端到端跑通 + 飞书 S1 卡片，第一条全链路。"),
            ("08:00–10:00 · DAY 2", "能力扩展 & 调优", "模型路由 / 企业审计 / 运行时纠偏 + Ultra（知识图谱 / What-if）。"),
            ("12:00 · DAY 2", "Demo & 提交", "真实 LLM 验证脚本 + 答辩录屏，可演示版本提交评审。"),
        ]
        for i, (t, h, p) in enumerate(timeline):
            x = Inches(0.65 + i * 3.05)
            add_rect(s7, x, Inches(1.65), Inches(2.9), Inches(2.8), PANEL, LINE)
            add_rect(s7, x, Inches(1.65), Inches(2.9), Inches(0.06), CYAN)
            add_textbox(s7, x + Inches(0.12), Inches(1.78), Inches(2.65), Inches(0.25), t, 9, CYAN, True)
            add_textbox(s7, x + Inches(0.12), Inches(2.1), Inches(2.65), Inches(0.35), h, 13, WHITE, True)
            add_textbox(s7, x + Inches(0.12), Inches(2.5), Inches(2.65), Inches(1.6), p, 10, DIM)
        add_rect(s7, Inches(0.65), Inches(4.7), Inches(12), Inches(0.75), RGBColor(0x33, 0x28, 0x00), GOLD)
        add_textbox(s7, Inches(0.8), Inches(4.82), Inches(11.5), Inches(0.55), "NEXT →  真实 Agent 平台事件回调双向集成：让 CoAgent 接管真实生产 Agent，从演示走向上线。", 12, DIM)
        add_footer(s7, 7)

        # Slide 8 — Q&A
        s8 = prs.slides.add_slide(blank)
        add_bg(s8)
        add_seclabel(s8, "Q & A", "欢迎提问")
        add_textbox(s8, Inches(0.65), Inches(1.5), Inches(6), Inches(0.9), "CoAgent", 56, CYAN, True)
        add_textbox(s8, Inches(0.65), Inches(2.45), Inches(8), Inches(0.7), "THANK YOU", 48, WHITE, True)
        add_textbox(s8, Inches(0.65), Inches(3.3), Inches(10), Inches(0.45), "为 Agent 上生产护航 —— 我们期待你的提问 · Beyond · 路睿 · 葛丽莎", 14, BLUE_L)
        add_textbox(s8, Inches(0.65), Inches(4.1), Inches(3), Inches(0.25), "QA 防翻车预案：", 10, MUTED, True)
        for i, q in enumerate([
            "「能自动修复吗？」→ 帮你处理 ≠ 替你处理：建议 + 人工审批 + 模拟执行，守住安全边界。",
            "「怎么接入现有 Agent？」→ 通过标准事件入口接入运行信号，按需适配业务工具 / MCP。",
            "「和 LangSmith / Datadog 区别？」→ 他们停在看见/诊断，我们做到一体化处置闭环。",
        ]):
            add_textbox(s8, Inches(0.65), Inches(4.4 + i * 0.42), Inches(11.5), Inches(0.4), "· " + q, 10, DIM)
        add_footer(s8, 8)

    return prs


def main() -> int:
    if not HTML.exists():
        print(f"Missing source: {HTML}", file=sys.stderr)
        return 1
    prs = build_presentation()
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Generated {OUT} ({len(prs.slides)} slides)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
